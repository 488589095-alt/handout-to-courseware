# -*- coding: utf-8 -*-
"""
build_pptx.py — 高中英语「讲义 → 课件」生成引擎（泛化版，不克隆任何标杆PPT）

输入：content.json(讲义抽取) + 原始课件模板.pptx [+ ai_scaffold.json(C类AI内容,可选)]
输出：完整课件 .pptx + slide_structure.json

两种讲次类型（--type 或 content.json 的 lecture_type 自动识别）：
  continuation  主观题·读后续写（第9讲型）：固化情节(基础/升级翻译) + 篇章(三步推理+情节描写)
  reading       客观题·阅读理解（第10讲型）：辅助方法 + 篇章训练(阅读每页1题·题→答)

设计原则（见 SKILL.md）：
  · A 讲义直出 / B 模版直出(封面/分隔/结束/品牌背景) / C AI生成(标"待师审")
  · 全部页面用"设计token"渲染，不克隆标杆；分页规则：阅读每页1题、翻译每页1题、完形每页3-4题
"""
import argparse, copy, json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
SHAPE_TAGS = ("sp", "pic", "grpSp", "graphicFrame", "cxnSp")
REF_ATTRS = [qn('r:embed'), qn('r:link'), qn('r:id'), qn('r:pict')]

# ════════ 设计 token（以 references/template_spec.md 为准，生成前必核对）════════
TITLE_COL = "BA7AC2"                     # 模版页标题色(spec: 课前/课后测 Text5)
PURPLE, PURPLE2 = "BA7AC2", "A076CE"     # 标题紫 / 次级紫(页码标pill)
INK, RED, GOLD = "1A1A1A", "FF5050", "F2A900"
CARD_BG, CARD_LN = "FFF7EC", "E6C9A8"    # 暖卡片
HL_BG, HL_LN = "E8F5E9", "46A35E"        # 正确高亮
EA_TITLE, EA_CN, EA_EN = "阿里巴巴普惠体 B", "可口可乐在乎体 楷体", "Arial"
PAGE_W, PAGE_H = 16.67, 12.5
CARD = (0.56, 1.04, 15.56, 10.77)        # 版式7 内容卡片区(spec)
AI_NOTE = "（AI 生成 · 待老师审核修改）"
CIRC = "①②③④⑤⑥⑦⑧⑨⑩"

# 模版 fixture 页（原始课件模板.pptx，1-based；spec 固定页）
T = dict(cover=2, titled=1, blank=4, div=[5, 7, 9], end=12)

SC = {}   # ai_scaffold（C类内容），render() 载入


# ════════ 基础：克隆/文本/表格 ════════
def clone_slide(prs, src):
    new = prs.slides.add_slide(src.slide_layout)
    spTree = new.shapes._spTree
    for ch in list(spTree):
        if etree.QName(ch).localname in SHAPE_TAGS:
            spTree.remove(ch)
    for el in src.shapes._spTree:
        if etree.QName(el).localname in SHAPE_TAGS:
            spTree.append(copy.deepcopy(el))
    relmap = {}
    for rId, rel in src.part.rels.items():
        if rel.reltype.endswith('slideLayout'):
            continue
        new_r = (new.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
                 if rel.is_external else new.part.relate_to(rel.target_part, rel.reltype))
        relmap[rId] = new_r
    for el in spTree.iter():
        for a in REF_ATTRS:
            if el.get(a) in relmap:
                el.set(a, relmap[el.get(a)])
    return new


def _style(run, size, bold, color, ea):
    """size/bold/color/ea 传 None = 不覆盖（保留模版继承）—— spec Gate 要求。"""
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = RGBColor.from_string(color)
    if ea is not None:
        rpr = run._r.get_or_add_rPr()
        for tag in ("ea", "latin"):
            e = rpr.find(f"{{{A_NS}}}{tag}")
            if e is None:
                e = etree.SubElement(rpr, f"{{{A_NS}}}{tag}")
            e.set("typeface", ea)


def tb(slide, text, left, top, w, h, size=24, bold=False, color=INK, ea=EA_EN,
       align=None, anchor=None, ls=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    first = True
    for line in str(text).split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        if align is not None:
            p.alignment = align
        if ls is not None:
            p.line_spacing = ls
        r = p.add_run(); r.text = line; _style(r, size, bold, color, ea)
    return box


def box(slide, text, left, top, w, h, fill=CARD_BG, line=CARD_LN, size=22,
        color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ea=EA_TITLE, bold=True):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(fill)
    sh.line.color.rgb = RGBColor.from_string(line); sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
    tf.clear(); first = True
    for ln in str(text).split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = align
        r = p.add_run(); r.text = ln; _style(r, size, bold, color, ea)
    return sh


def find(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh


def setname(slide, name, text, size=None, bold=None, ea=None):
    """换文字，样式以原 shape 显式属性为准；原本没写的属性不补写（继承模版）。"""
    sh = find(slide, name)
    if sh is None or not sh.has_text_frame:
        return
    r0 = next((r for p in sh.text_frame.paragraphs for r in p.runs), None)
    o_sz = r0.font.size.pt if (r0 and r0.font.size) else None
    o_bd = r0.font.bold if r0 else None
    o_cl = None
    try:
        if r0 and r0.font.color and r0.font.color.type is not None:
            o_cl = str(r0.font.color.rgb)
    except Exception:
        pass
    o_ea = None
    if r0 is not None:
        rpr = r0._r.find(f"{{{A_NS}}}rPr")
        if rpr is not None:
            e = rpr.find(f"{{{A_NS}}}ea")
            o_ea = e.get("typeface") if e is not None else None
    align0 = sh.text_frame.paragraphs[0].alignment
    tf = sh.text_frame; tf.clear(); first = True
    for ln in str(text).split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        if align0 is not None:
            p.alignment = align0
        r = p.add_run(); r.text = ln
        _style(r, size or o_sz, bold if bold is not None else o_bd,
               o_cl, ea or o_ea)


def add_table(slide, rows, left, top, w, h, col_w, head=True, hsz=18, bsz=16):
    nr, nc = len(rows), len(rows[0])
    gf = slide.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(w), Inches(h))
    t = gf.table
    for c, cw in enumerate(col_w):
        t.columns[c].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, txt in enumerate(row):
            cell = t.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.1)
            tf = cell.text_frame; tf.word_wrap = True; tf.clear()
            p = tf.paragraphs[0]
            is_head = head and ri == 0
            p.alignment = PP_ALIGN.CENTER if (is_head or nc > 2) else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(txt)
            _style(r, hsz if is_head else bsz, is_head, INK,
                   EA_TITLE if is_head else "微软雅黑")
    return gf


def _blank(prs):
    """模版空白内容页（品牌背景来自模版版式）。"""
    s = clone_slide(prs, prs.slides[T["blank"] - 1])
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def _titled(prs, title):
    """带标题内容页（Gate：标题样式必须来自模版，不得自创）：
       克隆模版第1页(版式7=卡片框+标题槽 Text 5)，仅换标题文字 →
       自动继承 spec 样式 sz40/不加粗/#BA7AC2/阿里巴巴普惠体B/居中@y0.98。"""
    s = clone_slide(prs, prs.slides[T["titled"] - 1])
    setname(s, "Text 5", title)
    sh = find(s, "Text 5")                 # 长标题加宽，保持居中
    if sh is not None:
        sh.left = Inches((PAGE_W - 12.0) / 2); sh.width = Inches(12.0)
    return s


def ai_note(s):
    tb(s, AI_NOTE, 1.0, PAGE_H - 1.0, PAGE_W - 2, 0.5, size=14, color=PURPLE2,
       ea="微软雅黑", align=PP_ALIGN.CENTER)


# ════════ 通用页（两种讲次共用）════════
def r_cover(prs, C):
    s = clone_slide(prs, prs.slides[T["cover"] - 1])
    setname(s, "Text 0", C["lecture_no"])
    setname(s, "Text 1", C["title"])
    return s


def r_toc(prs, C):
    s = _titled(prs, "目录")
    y = 3.2
    for i, p in enumerate(C["parts"], 1):
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.2), Inches(y), Inches(1.0), Inches(1.0))
        c.fill.solid(); c.fill.fore_color.rgb = RGBColor.from_string(PURPLE)
        c.line.fill.background(); c.shadow.inherit = False
        tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.clear()
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run(); rr.text = str(i); _style(rr, 30, True, "FFFFFF", EA_TITLE)
        tb(s, f"{p['part_label']}   {p['title']}", 4.6, y, 10.5, 1.0, size=32, bold=True,
           color=INK, ea=EA_TITLE, anchor=MSO_ANCHOR.MIDDLE)
        y += 1.8
    return s


def r_km(prs, img_path):
    s = _titled(prs, "Knowledge Map")
    if img_path and Path(img_path).exists():
        s.shapes.add_picture(str(img_path), Inches((PAGE_W - 13) / 2), Inches(2.6),
                             width=Inches(13))
    return s


def r_preview(prs, C):
    s = _titled(prs, "Preview")
    pv = C["preview"]
    add_table(s, [pv["header"]] + pv["rows"], 1.2, 2.4, 14.2, 5.5,
              [1.7, 3.0, 2.3, 7.2], hsz=20, bsz=18)
    return s


def r_leadin(prs, C):
    s = _titled(prs, "Leading-in")
    tb(s, C.get("leading_in", ""), 1.8, 4.5, 13.0, 4.0, size=30, color=INK,
       ea=EA_CN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ls=1.4)
    return s


def r_divider(prs, part, no):
    s = clone_slide(prs, prs.slides[T["div"][min(no, len(T["div"])) - 1] - 1])
    setname(s, "Text 1", part["part_label"])
    setname(s, "Text 0", part.get("subtitle") or part["title"])
    return s


def r_section(prs, title):
    s = _blank(prs)
    tb(s, title, 1.0, 4.6, PAGE_W - 2, 1.8, size=72, bold=True, color=PURPLE,
       ea=EA_CN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


def r_source(prs, src, level=""):
    s = _blank(prs)
    tb(s, str(src).replace("【", "").replace("】", ""), 1.0, 5.0, PAGE_W - 2, 1.6,
       size=60, bold=True, color=INK, ea=EA_CN, align=PP_ALIGN.CENTER,
       anchor=MSO_ANCHOR.MIDDLE)
    if level:
        tb(s, level, 1.0, 6.8, PAGE_W - 2, 0.8, size=28, color=GOLD, ea=EA_EN,
           align=PP_ALIGN.CENTER)
    return s


def r_passage(prs, paras, tail=""):
    s = _blank(prs)
    body = "\n".join(paras) + (("\n\n" + tail) if tail else "")
    tb(s, body, 0.9, 0.9, PAGE_W - 1.8, PAGE_H - 1.7, size=18, color=INK, ea=EA_EN, ls=1.15)
    return s


def r_end(prs, C):
    s = clone_slide(prs, prs.slides[T["end"] - 1])
    setname(s, "Text 0", C["end"]["big"])
    setname(s, "Text 1", C["end"]["small"])
    return s


# ════════ reading（客观题·阅读）════════
def r_table_page(prs, title, rows, col_w, head=True, bsz=16):
    s = _titled(prs, title)
    add_table(s, rows, 1.0, 2.4, PAGE_W - 2, 8.0, col_w, head=head, bsz=bsz)
    return s


def r_question(prs, tag, stem, options):
    s = _titled(prs, tag)
    tb(s, stem + "\n" + "\n".join(options), 1.0, 2.2, PAGE_W - 2, PAGE_H - 3.0,
       size=26, color=INK, ea=EA_EN, ls=1.3)
    return s


def r_answer(prs, tag, answer, analysis):
    s = _titled(prs, tag)
    tb(s, f"【答案】{answer}", 1.0, 2.2, PAGE_W - 2, 0.9, size=30, bold=True,
       color=RED, ea=EA_CN)
    tb(s, "【解析】" + analysis, 1.0, 3.2, PAGE_W - 2, PAGE_H - 4.0, size=20,
       color=INK, ea=EA_CN, ls=1.25)
    return s


# ════════ continuation（主观题·读后续写）════════
def r_flow(prs, title, steps):
    """固化情节5步概览（C·AI 或 讲义情节名）。"""
    s = _titled(prs, title)
    n = len(steps); bw = 2.6; gap = (PAGE_W - n * bw) / (n + 1)
    for i, name in enumerate(steps):
        x = gap + i * (bw + gap)
        box(s, f"{i+1}\n{name}", x, 5.0, bw, 2.3, fill="F3E9F7", line=PURPLE2,
            size=26, color=PURPLE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + bw + 0.05),
                                    Inches(5.9), Inches(gap - 0.1), Inches(0.5))
            ar.fill.solid(); ar.fill.fore_color.rgb = RGBColor.from_string("CBA6E0")
            ar.line.fill.background(); ar.shadow.inherit = False
    return s


def r_guide(prs, name, prompt="你会怎么表达?"):
    s = _blank(prs)
    box(s, name, 6.6, 3.2, 3.4, 2.2, fill="F3E9F7", line=PURPLE2, size=44,
        color=PURPLE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ea=EA_CN)
    tb(s, prompt, 1.0, 7.4, PAGE_W - 2, 1.8, size=72, color=INK, ea=EA_CN,
       align=PP_ALIGN.CENTER)
    return s


def r_translation(prs, label, name, zh, en, note=""):
    """翻译页（基础版本/升级版本）：中文 + 英文（讲义直出）。"""
    s = _blank(prs)
    box(s, label, 0.99, 0.98, 3.46, 1.2, fill="F3E9F7", line=PURPLE2, size=40,
        color=PURPLE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    box(s, name, 13.0, 0.95, 2.6, 2.2, fill="FFFFFF", line=PURPLE2, size=40,
        color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ea=EA_CN)
    tb(s, zh, 1.27, 5.0, 14.14, 1.6, size=34, color=INK, ea=EA_CN, ls=1.2)
    tb(s, en, 1.27, 6.9, 14.14, 1.6, size=27, color=INK, ea=EA_EN, ls=1.15)
    if note:
        tb(s, note, 1.0, PAGE_H - 1.0, PAGE_W - 2, 0.5, size=13, color=PURPLE2,
           ea="微软雅黑")
    return s


def r_step(prs, steps, active):
    s = _blank(prs)
    tb(s, "阅读 · 解题三步", 1.0, 2.0, PAGE_W - 2, 1.0, size=42, bold=True,
       color=PURPLE, ea=EA_TITLE, align=PP_ALIGN.CENTER)
    n = len(steps); bw = 4.2; gap = (PAGE_W - n * bw) / (n + 1)
    for i, st in enumerate(steps):
        on = i == active
        box(s, st, gap + i * (bw + gap), 5.5, bw, 2.3,
            fill=PURPLE if on else "F3E9F7", line=PURPLE if on else "CBA6E0",
            size=30, color="FFFFFF" if on else "B49AC9",
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


def r_bilingual(prs, en_paras, zh_paras):
    s = _titled(prs, "全文翻译对照")
    tb(s, "\n".join(en_paras), 0.8, 2.2, 7.55, 9.3, size=16, color=INK, ea=EA_EN, ls=1.12)
    tb(s, "\n".join(zh_paras), 8.45, 2.2, 7.55, 9.3, size=17, color=INK, ea=EA_CN, ls=1.18)
    ai_note(s)
    return s


def r_overall(prs, summary, direction, reveal):
    """整体走向：原文概括 →(揭示) 后续走向，两页式。C·AI。"""
    s = _titled(prs, "情节推理 · 整体走向")
    box(s, "原文概括\n" + summary, 1.4, 2.2, 13.8, 3.3)
    ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(7.6), Inches(5.7),
                            Inches(1.4), Inches(1.0))
    ar.fill.solid(); ar.fill.fore_color.rgb = RGBColor.from_string(PURPLE2)
    ar.line.fill.background(); ar.shadow.inherit = False
    if reveal:
        box(s, "后续走向（整体方向）\n" + direction, 1.4, 6.9, 13.8, 3.4,
            fill="F3E9F7", line=PURPLE2)
    ai_note(s)
    return s


def r_reason(prs, q, reveal):
    """细节推理：原文细节 + 推测选择题，题/答两页式。C·AI。"""
    s = _titled(prs, "情节推理 · 细节推理")
    box(s, q.get("context", ""), 1.4, 2.1, 13.8, 1.7, anchor=MSO_ANCHOR.MIDDLE)
    tb(s, "Q：" + q.get("stem", ""), 1.4, 4.0, 13.8, 0.7, size=26, bold=True,
       color="7D292D", ea=EA_CN)
    y = 4.9
    for i, opt in enumerate(q.get("options", [])):
        on = reveal and i == q.get("answer", 0)
        box(s, f"{chr(65+i)}.  {opt}", 1.8, y, 13.0, 1.5,
            fill=HL_BG if on else "FFFFFF", line=HL_LN if on else "CBA6E0",
            size=22, color=INK, anchor=MSO_ANCHOR.MIDDLE, ea=EA_CN, bold=False)
        y += 1.75
    if reveal:
        tb(s, f"【答案】{chr(65 + q.get('answer', 0))}", 1.4, y + 0.1, 13.8, 0.7,
           size=26, bold=True, color=RED, ea=EA_CN)
    ai_note(s)
    return s


def r_segment(prs, details, filled, hint1="", hint2=""):
    """情节切分总表（首段①–④ / 次段⑤–⑧），留白→填好两页式。内容=讲义。"""
    s = _titled(prs, "情节切分" + ("" if filled else "（留白练习）"))
    def rows(ds):
        return [[CIRC[d["index"] - 1],
                 (f"【{d['label']}】{d['zh']}" if filled else "")] for d in ds]
    shou = [d for d in details if d["group"] == "首段"]
    ci = [d for d in details if d["group"] == "次段"]
    tb(s, "首段" + (f"提示句：{hint1}" if hint1 else ""), 1.2, 2.1, 14.2, 0.5,
       size=22, bold=True, color="7D292D", ea=EA_TITLE)
    add_table(s, rows(shou), 1.2, 2.65, 14.2, 3.4, [0.9, 13.3], head=False, bsz=15)
    tb(s, "次段" + (f"提示句：{hint2}" if hint2 else ""), 1.2, 6.45, 14.2, 0.5,
       size=22, bold=True, color="7D292D", ea=EA_TITLE)
    add_table(s, rows(ci), 1.2, 7.0, 14.2, 3.4, [0.9, 13.3], head=False, bsz=15)
    return s


def r_detail(prs, d, desc, reveal):
    """情节描写：中文 →(揭示) 英文，左上 序号+情节名+首/次段+描写类型(C·AI)。"""
    s = _blank(prs)
    box(s, d["group"], 0.94, 1.07, 1.4, 1.65, fill="F3E9F7", line=PURPLE2, size=34,
        color=PURPLE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ea=EA_CN)
    box(s, f"{CIRC[d['index']-1]} {d['label']}", 2.45, 1.07, 3.11, 0.79,
        fill="FFFFFF", line=PURPLE2, size=30, color=INK,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ea=EA_CN)
    box(s, desc, 2.45, 1.93, 3.11, 0.79, fill="FFFFFF", line=PURPLE2, size=30,
        color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, ea=EA_CN)
    tb(s, d["zh"], 1.29, 5.1, 14.09, 1.6, size=30, color=INK, ea=EA_CN, ls=1.2)
    if reveal:
        tb(s, d["en"], 1.29, 6.9, 14.09, 1.6, size=26, color=INK, ea=EA_EN, ls=1.15)
    return s


def r_essay(prs, part):
    s = _blank(prs)
    box(s, "参考范文", 0.9, 1.07, 2.47, 0.79, fill="F3E9F7", line=PURPLE2, size=30,
        color=PURPLE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    d = part["details"]
    p1 = part["prompts"]["para1"].replace("Para 1: ", "段落一：") + " " + \
        " ".join(x["en"] for x in d if x["group"] == "首段")
    p2 = part["prompts"]["para2"].replace("Para 2: ", "段落二：") + " " + \
        " ".join(x["en"] for x in d if x["group"] == "次段")
    tb(s, p1 + "\n\n" + p2, 1.29, 2.3, 14.09, 9.3, size=24, bold=True, color=INK,
       ea=EA_EN, ls=1.2)
    return s


# ════════ 编排 ════════
def specs_reading(C):
    sp = []
    add = lambda role, summary, **kw: sp.append({"role": role, "summary": summary, **kw})
    add("cover", "封面"); add("toc", "目录"); add("km", "Knowledge Map")
    add("preview", "Preview"); add("leadin", "Leading-in")
    add("divider", "PART1", part=C["parts"][0], no=1)
    zz = C["part1_zhuzhi"]
    add("section", zz.get("title", "主旨辅助"))
    add("source", f"例篇 {zz['source']}", src=zz["source"])
    add("passage", "例篇语篇", paras=zz["passage"])
    if zz.get("method_table"):
        add("table", "主旨辅助 练习表", rows=zz["method_table"], col_w=[5.0, 9.67], head=False)
    for i, q in enumerate(zz["questions"], 1):
        tag = f"主旨辅助 · 第{i}题"
        add("question", tag, tag=tag, stem=q["stem"], options=q["options"])
        add("answer", tag + " 答", tag=tag, answer=q.get("answer", ""),
            analysis=q.get("analysis", ""))
    xx = C["part1_xuanxiang"]
    add("section", xx.get("title", "选项辅助"))
    if xx.get("table"):
        add("table", "选项辅助 特征表", rows=xx["table"], col_w=[2.6, 4.6, 2.6, 4.87],
            head=True, bsz=14)
    add("divider", "PART2", part=C["parts"][1], no=2)
    for psg in C["part2"]:
        add("source", f"{psg['name']} {psg['source']}", src=psg["source"],
            level=psg.get("level", ""))
        add("passage", f"{psg['name']} 语篇", paras=psg["passage"])
        for i, q in enumerate(psg["questions"], 1):
            tag = f"{psg['name']} · 第{i}题"
            add("question", tag, tag=tag, stem=q["stem"], options=q["options"])
            add("answer", tag + " 答", tag=tag, answer=q.get("answer", ""),
                analysis=q.get("analysis", ""))
    add("end", "结束")
    return sp


def specs_continuation(C):
    sp = []
    add = lambda role, summary, **kw: sp.append({"role": role, "summary": summary, **kw})
    add("cover", "封面"); add("toc", "目录"); add("km", "Knowledge Map")
    add("preview", "Preview")
    add("divider", "PART1", part=C["parts"][0], no=1)
    add("leadin", "Leading-in")
    add("flow", "固化情节概览",
        title=SC.get("plot_overview", {}).get("title", "固化情节"),
        steps=SC.get("plot_overview", {}).get("steps",
              [p["name"] for p in C["part1_plots"]]))
    for plot in C["part1_plots"]:
        add("guide", f"引导·{plot['name']}", name=plot["name"])
        add("trans", f"{plot['name']} 基础版本", label="基础版本", name=plot["name"],
            zh=plot["basic"]["zh"], en=plot["basic"]["en"])
        add("trans", f"{plot['name']} 升级版本", label="升级版本", name=plot["name"],
            zh=plot["upgrade"]["zh"], en=plot["upgrade"]["en"])
    steps = SC.get("step_framework", ["步骤一·原文概括", "步骤二·细节推理", "步骤三·段首分析"])
    for pi, key in [(1, "part2"), (2, "part3")]:
        part = C[key]; sc = SC.get(key, {})
        add("divider", f"PART{pi+1}", part=C["parts"][pi], no=pi + 1)
        add("source", f"题源 {part['source']}", src=part["source"])
        p1 = part["prompts"]["para1"].replace("Para 1: ", "段落一：")
        p2 = part["prompts"]["para2"].replace("Para 2: ", "段落二：")
        add("passage", "语篇+续写提示", paras=part["passage_en"], tail=p1 + "\n" + p2)
        add("step", "步骤一", steps=steps, active=0)
        if part.get("passage_zh"):
            add("bilingual", "全文对照(C·AI翻译)", en=part["passage_en"], zh=part["passage_zh"])
        add("overall", "整体走向·概括", osum=sc.get("overall_summary", ""),
            direction=sc.get("overall_direction", ""), reveal=False)
        add("overall", "整体走向·揭示", osum=sc.get("overall_summary", ""),
            direction=sc.get("overall_direction", ""), reveal=True)
        add("step", "步骤二", steps=steps, active=1)
        if sc.get("reasoning_q"):
            add("reason", "细节推理·题", q=sc["reasoning_q"], reveal=False)
            add("reason", "细节推理·答", q=sc["reasoning_q"], reveal=True)
        add("step", "步骤三", steps=steps, active=2)
        h1 = part["prompts"]["para1"].replace("Para 1: ", "")
        h2 = part["prompts"]["para2"].replace("Para 2: ", "")
        add("segment", "情节切分·留白", details=part["details"], filled=False, h1=h1, h2=h2)
        add("segment", "情节切分·填好", details=part["details"], filled=True, h1=h1, h2=h2)
        add("section", "情节描写")
        descs = sc.get("desc_types", [])
        for d in part["details"]:
            desc = descs[d["index"] - 1] if d["index"] - 1 < len(descs) else "情节描写"
            tagn = f"{CIRC[d['index']-1]}{d['label']}"
            add("detail", f"情节描写 {tagn} 题", d=d, desc=desc, reveal=False)
            add("detail", f"情节描写 {tagn} 答", d=d, desc=desc, reveal=True)
        add("essay", "参考范文", part=part)
    add("end", "结束")
    return sp




# ════════ Gate：生成前核对模版拆解文档 ════════
def check_template(prs, content_path):
    """生成前必过 Gate：模版必须与 references/template_spec.json 一致（字体/字号/标题槽/fixture）。"""
    cand = [Path(content_path).parent / "template_spec.json",
            Path(__file__).resolve().parent.parent / "references" / "template_spec.json"]
    spec_path = next((p for p in cand if p.exists()), None)
    if spec_path is None:
        raise SystemExit("❌ Gate未过：缺 template_spec.json，先运行 dissect_template.py 拆解模版")
    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    errs = []
    need = max(T["cover"], T["titled"], T["blank"], T["end"], max(T["div"]))
    if len(prs.slides) < need:
        errs.append(f"模版仅{len(prs.slides)}页 < fixture索引{need}")
    s1 = next((s for s in spec["slides"] if s["no"] == T["titled"]), {})
    t5 = next((sh for sh in s1.get("shapes", []) if sh.get("name") == "Text 5"), None)
    if not t5:
        errs.append("spec第1页缺标题槽 Text 5")
    else:
        f = t5.get("font", {})
        if f.get("sz") != 40.0 or f.get("color") != TITLE_COL or \
                f.get("ea") != EA_TITLE:
            errs.append(f"标题槽样式漂移: spec={f} ≠ token(sz40/{TITLE_COL}/{EA_TITLE})")
    cov = next((s for s in spec["slides"] if s["no"] == T["cover"]), {})
    if not any(sh.get("name") == "Text 0" for sh in cov.get("shapes", [])):
        errs.append("spec封面页缺 Text 0(讲次槽)")
    if errs:
        print("❌ 模版核对未通过（修 token 或重跑 dissect_template.py）:")
        for e in errs:
            print("   -", e)
        raise SystemExit(1)
    print(f"✅ Gate通过：模版与 {spec_path.name} 一致")


def add_sections(prs, specs):
    """PowerPoint 节（大纲分组）：封面·目录 / 各PART / 结束。"""
    P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
    SEC = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"
    GUIDS = [f"{{B1A9C8E0-000{i}-4A1A-9C01-00000000000{i}}}" for i in range(1, 9)]
    bounds = [("封面 · 目录", 1)]
    for i, s in enumerate(specs):
        if s["role"] == "divider":
            bounds.append((s["summary"], i + 1))
    end_pos = next((i + 1 for i, s in enumerate(specs) if s["role"] == "end"), None)
    if end_pos:
        bounds.append(("结束", end_pos))
    sld_ids = [el.get("id") for el in prs.slides._sldIdLst]
    pres = prs.part._element
    extLst = pres.find(qn("p:extLst"))
    if extLst is None:
        extLst = etree.SubElement(pres, qn("p:extLst"))
    for ex in list(extLst.findall(qn("p:ext"))):
        if ex.get("uri") == SEC:
            extLst.remove(ex)
    ext = etree.SubElement(extLst, qn("p:ext")); ext.set("uri", SEC)
    secLst = etree.SubElement(ext, f"{{{P14}}}sectionLst")
    for si, (name, start) in enumerate(bounds):
        end = bounds[si + 1][1] - 1 if si + 1 < len(bounds) else len(sld_ids)
        sec = etree.SubElement(secLst, f"{{{P14}}}section")
        sec.set("name", name); sec.set("id", GUIDS[si])
        lst = etree.SubElement(sec, f"{{{P14}}}sldIdLst")
        for idx in range(start, end + 1):
            sid = etree.SubElement(lst, f"{{{P14}}}sldId")
            sid.set("id", sld_ids[idx - 1])


def render(content_path, template_path, out_path, lec_type=None, km_image=None):
    global SC
    C = json.loads(Path(content_path).read_text(encoding="utf-8"))
    scp = Path(content_path).with_name("ai_scaffold.json")
    SC = json.loads(scp.read_text(encoding="utf-8")) if scp.exists() else {}
    lec_type = lec_type or C.get("lecture_type") or \
        ("continuation" if "part1_plots" in C else "reading")
    specs = specs_continuation(C) if lec_type == "continuation" else specs_reading(C)

    Path(out_path).with_name("slide_structure.json").write_text(
        json.dumps({"type": lec_type, "total": len(specs),
                    "slides": [{"n": i + 1, "role": s["role"], "summary": s["summary"]}
                               for i, s in enumerate(specs)]},
                   ensure_ascii=False, indent=2), encoding="utf-8")

    km = km_image or str(Path(content_path).parent / "assets" / "knowledge_map.png")
    prs = Presentation(str(template_path))
    check_template(prs, content_path)          # ← Gate：生成前核对模版拆解文档
    orig = len(prs.slides)
    for s in specs:
        r = s["role"]
        if r == "cover": r_cover(prs, C)
        elif r == "toc": r_toc(prs, C)
        elif r == "km": r_km(prs, km)
        elif r == "preview": r_preview(prs, C)
        elif r == "leadin": r_leadin(prs, C)
        elif r == "divider": r_divider(prs, s["part"], s["no"])
        elif r == "section": r_section(prs, s["summary"])
        elif r == "source": r_source(prs, s["src"], s.get("level", ""))
        elif r == "passage": r_passage(prs, s["paras"], s.get("tail", ""))
        elif r == "table": r_table_page(prs, s["summary"], s["rows"], s["col_w"],
                                        s.get("head", True), s.get("bsz", 16))
        elif r == "question": r_question(prs, s["tag"], s["stem"], s["options"])
        elif r == "answer": r_answer(prs, s["tag"], s["answer"], s["analysis"])
        elif r == "flow": r_flow(prs, s["title"], s["steps"])
        elif r == "guide": r_guide(prs, s["name"])
        elif r == "trans": r_translation(prs, s["label"], s["name"], s["zh"], s["en"])
        elif r == "step": r_step(prs, s["steps"], s["active"])
        elif r == "bilingual": r_bilingual(prs, s["en"], s["zh"])
        elif r == "overall": r_overall(prs, s["osum"], s["direction"], s["reveal"])
        elif r == "reason": r_reason(prs, s["q"], s["reveal"])
        elif r == "segment": r_segment(prs, s["details"], s["filled"],
                                       s.get("h1", ""), s.get("h2", ""))
        elif r == "detail": r_detail(prs, s["d"], s["desc"], s["reveal"])
        elif r == "essay": r_essay(prs, s["part"])
        elif r == "end": r_end(prs, C)
    for i in range(orig, 0, -1):
        el = list(prs.slides._sldIdLst)[i - 1]
        prs.part.drop_rel(el.get(qn("r:id"))); prs.slides._sldIdLst.remove(el)
    try:
        add_sections(prs, specs)                # 大纲(节)分组
    except Exception as e:
        print(f"⚠️ 节添加失败(忽略): {e}")
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"✅ 课件已生成（{lec_type}）：{out_path}（{len(prs.slides)} 页）")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--content", required=True)
    ap.add_argument("--template", required=True)
    ap.add_argument("--type", choices=["continuation", "reading"], default=None)
    ap.add_argument("--km-image", default=None)
    ap.add_argument("-o", "--output", required=True)
    a = ap.parse_args()
    render(a.content, a.template, a.output, a.type, a.km_image)
