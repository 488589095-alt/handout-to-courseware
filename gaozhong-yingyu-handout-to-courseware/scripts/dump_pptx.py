# -*- coding: utf-8 -*-
"""Dump pptx structure (per-slide shapes/text/pos) and docx body (paragraphs+tables)
   to plain text for analysis. Usage:
     python3 dump_files.py pptx <file.pptx> -o out.txt
     python3 dump_files.py docx <file.docx> -o out.txt
"""
import argparse, sys
from pptx import Presentation
from pptx.util import Emu

def emu_in(v):
    try:
        return round(Emu(v).inches, 2)
    except Exception:
        return None

def dump_pptx(path, out):
    prs = Presentation(path)
    lines = []
    lines.append(f"# {path}")
    lines.append(f"slide_size = {emu_in(prs.slide_width)} x {emu_in(prs.slide_height)} in")
    lines.append(f"total_slides = {len(prs.slides)}")
    for i, slide in enumerate(prs.slides, 1):
        try:
            lay = slide.slide_layout.name
        except Exception:
            lay = "?"
        lines.append(f"\n===== SLIDE {i} =====")
        lines.append(f"[layout: {lay}]")
        for sh in slide.shapes:
            try:
                st = sh.shape_type
            except Exception:
                st = "?"
            pos = f"pos=({emu_in(sh.left)},{emu_in(sh.top)})" if sh.left is not None else "pos=?"
            size = f"size=({emu_in(sh.width)}x{emu_in(sh.height)})" if sh.width is not None else "size=?"
            tag = ""
            if sh.shape_type is not None and "PICTURE" in str(sh.shape_type):
                tag = "[PICTURE]"
            elif sh.has_table:
                tag = "[TABLE]"
            lines.append(f"  <{st}> name='{sh.name}' {pos} {size} {tag}")
            if sh.has_text_frame and sh.text_frame.text.strip():
                # collapse paragraphs with ' / '
                txt = " / ".join(p.text for p in sh.text_frame.text_frame.paragraphs) if False else \
                      " / ".join(p.text for p in sh.text_frame.paragraphs)
                # show run-level font info for first non-empty run
                fontinfo = []
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.text.strip():
                            sz = r.font.size.pt if r.font.size else None
                            col = None
                            try:
                                col = str(r.font.color.rgb) if r.font.color and r.font.color.type is not None else None
                            except Exception:
                                col = None
                            ea = None
                            try:
                                rpr = r._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                            except Exception:
                                rpr = None
                            fontinfo.append(f"sz={sz},bold={r.font.bold},col={col},name={r.font.name}")
                            break
                    if fontinfo:
                        break
                lines.append(f"      TEXT: {txt}")
                if fontinfo:
                    lines.append(f"      FONT: {fontinfo[0]}")
            if sh.has_table:
                tbl = sh.table
                lines.append(f"      TABLE {len(tbl.rows)}x{len(tbl.columns)}:")
                for r in tbl.rows:
                    cells = [c.text.replace('\n', '⏎') for c in r.cells]
                    lines.append("        | " + " | ".join(cells))
    open(out, "w", encoding="utf-8").write("\n".join(lines))
    print(f"wrote {out} ({len(lines)} lines)")

def dump_docx(path, out):
    import docx
    from docx.document import Document as _Doc
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    doc = docx.Document(path)
    lines = [f"# {path}"]
    body = doc.element.body
    pi = ti = 0
    for child in body.iterchildren():
        if isinstance(child, CT_P):
            p = Paragraph(child, doc)
            style = p.style.name if p.style else ""
            txt = p.text
            if txt.strip():
                # show bold/size of first run
                meta = ""
                for r in p.runs:
                    if r.text.strip():
                        meta = f"[b={r.bold},sz={r.font.size.pt if r.font.size else None},name={r.font.name}]"
                        break
                lines.append(f"P{pi} <{style}> {meta}: {txt}")
            pi += 1
        elif isinstance(child, CT_Tbl):
            t = Table(child, doc)
            lines.append(f"\n--- TABLE {ti} ({len(t.rows)}x{len(t.columns)}) ---")
            for row in t.rows:
                cells = [c.paragraphs[0].text if c.paragraphs else "" for c in row.cells]
                lines.append("  | " + " | ".join(cells))
            lines.append("")
            ti += 1
    open(out, "w", encoding="utf-8").write("\n".join(lines))
    print(f"wrote {out} ({len(lines)} lines, {pi} paras, {ti} tables)")

if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("kind", choices=["pptx", "docx"])
    ap.add_argument("file")
    ap.add_argument("-o", "--out", required=True)
    a = ap.parse_args()
    if a.kind == "pptx":
        dump_pptx(a.file, a.out)
    else:
        dump_docx(a.file, a.out)
