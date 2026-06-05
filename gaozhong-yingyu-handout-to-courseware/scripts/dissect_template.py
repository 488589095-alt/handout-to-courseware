# -*- coding: utf-8 -*-
"""
dissect_template.py — 把 PPT 模版拆解成可完全理解的规格文档（生成前必读/必核对的 Gate）

输出：
  template_spec.md   人读：主题字体/配色、母版与版式占位符、每张固定页的元素(坐标/字体/字号/颜色)
  template_spec.json 机读：build 前自动核对（字体、字号、标题位置等 token 必须以此为准）

用法: python3 dissect_template.py <模版.pptx> -o <outdir>
"""
import argparse, json
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def IN(v):
    try:
        return round(Emu(v).inches, 2)
    except Exception:
        return None


def run_props(r):
    p = {}
    if r.font.size:
        p["sz"] = r.font.size.pt
    if r.font.bold is not None:
        p["bold"] = r.font.bold
    try:
        if r.font.color and r.font.color.type is not None:
            p["color"] = str(r.font.color.rgb)
    except Exception:
        pass
    rpr = r._r.find(f"{A}rPr")
    if rpr is not None:
        for tag in ("latin", "ea"):
            e = rpr.find(f"{A}{tag}")
            if e is not None:
                p[tag] = e.get("typeface")
    return p


def shape_info(sh):
    d = {"name": sh.name, "type": str(sh.shape_type)}
    if sh.left is not None:
        d["pos"] = [IN(sh.left), IN(sh.top)]
        d["size"] = [IN(sh.width), IN(sh.height)]
    if getattr(sh, "is_placeholder", False):
        ph = sh.placeholder_format
        d["ph"] = {"idx": ph.idx, "type": str(ph.type)}
    if sh.has_text_frame:
        tf = sh.text_frame
        if tf.text.strip():
            d["text"] = " / ".join(p.text for p in tf.paragraphs)[:120]
        for p in tf.paragraphs:
            for r in p.runs:
                rp = run_props(r)
                if rp:
                    d["font"] = rp
                    if p.alignment is not None:
                        d["align"] = str(p.alignment)
                    break
            if "font" in d:
                break
    return d


def theme_info(prs):
    th = prs.slide_masters[0].element.getroottree().getroot()  # not theme; use part
    # theme via master part rels
    master = prs.slide_masters[0]
    theme = None
    for rel in master.part.rels.values():
        if "theme" in rel.reltype:
            theme = rel.target_part
            break
    out = {"major_font": None, "minor_font": None, "colors": {}}
    if theme is None:
        return out
    root = theme._element if hasattr(theme, "_element") else None
    import lxml.etree as ET
    root = ET.fromstring(theme.blob)
    scheme = root.find(f".//{A}fontScheme")
    if scheme is not None:
        for tag, key in (("majorFont", "major_font"), ("minorFont", "minor_font")):
            f = scheme.find(f"{A}{tag}")
            if f is not None:
                lat = f.find(f"{A}latin")
                ea = f.find(f"{A}ea")
                out[key] = {"latin": lat.get("typeface") if lat is not None else None,
                            "ea": ea.get("typeface") if ea is not None else None}
    clr = root.find(f".//{A}clrScheme")
    if clr is not None:
        for c in clr:
            tag = c.tag.split("}")[1]
            srgb = c.find(f"{A}srgbClr")
            sysc = c.find(f"{A}sysClr")
            out["colors"][tag] = (srgb.get("val") if srgb is not None
                                  else (sysc.get("lastClr") if sysc is not None else None))
    return out


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("template")
    ap.add_argument("-o", "--outdir", required=True)
    a = ap.parse_args()
    prs = Presentation(a.template)
    spec = {"file": Path(a.template).name,
            "page_in": [IN(prs.slide_width), IN(prs.slide_height)],
            "theme": theme_info(prs), "layouts": [], "slides": []}

    for m in prs.slide_masters:
        for L in m.slide_layouts:
            spec["layouts"].append({"name": L.name,
                                    "shapes": [shape_info(sh) for sh in L.shapes]})
    for i, s in enumerate(prs.slides, 1):
        spec["slides"].append({"no": i, "layout": s.slide_layout.name,
                               "shapes": [shape_info(sh) for sh in s.shapes]})

    out = Path(a.outdir); out.mkdir(parents=True, exist_ok=True)
    (out / "template_spec.json").write_text(json.dumps(spec, ensure_ascii=False, indent=1),
                                            encoding="utf-8")

    md = [f"# 模版拆解：{spec['file']}", "",
          f"- 画布：{spec['page_in'][0]} × {spec['page_in'][1]} in",
          f"- 主题字体：major={spec['theme']['major_font']}  minor={spec['theme']['minor_font']}",
          f"- 主题色：{spec['theme']['colors']}", "", "## 版式（layouts）"]
    for L in spec["layouts"]:
        md.append(f"\n### {L['name']}")
        for sh in L["shapes"]:
            md.append(f"- `{sh['name']}` {sh.get('type','')} pos={sh.get('pos')} size={sh.get('size')}"
                      + (f" ph={sh['ph']}" if 'ph' in sh else "")
                      + (f" font={sh.get('font')}" if sh.get('font') else "")
                      + (f"  text=「{sh.get('text','')[:40]}」" if sh.get('text') else ""))
    md.append("\n## 固定页（slides）")
    for s in spec["slides"]:
        md.append(f"\n### 第{s['no']}页  [layout: {s['layout']}]")
        for sh in s["shapes"]:
            md.append(f"- `{sh['name']}` pos={sh.get('pos')} size={sh.get('size')}"
                      + (f" font={sh.get('font')}" if sh.get('font') else "")
                      + (f" align={sh.get('align')}" if sh.get('align') else "")
                      + (f"  text=「{sh.get('text','')[:50]}」" if sh.get('text') else ""))
    (out / "template_spec.md").write_text("\n".join(md), encoding="utf-8")
    print(f"✅ template_spec.md / template_spec.json → {out}")


if __name__ == "__main__":
    main()
