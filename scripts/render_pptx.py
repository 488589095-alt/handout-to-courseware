"""
render_pptx.py v0.1 — 基于"修改模板已有 sample slide"模式渲染

策略（v0 add_slide 跑版的修复方案）：
  - 不删除模板已有 23 张 sample slide（这些含装饰图，必须保留）
  - 按 slide_structure.json 指定的 source 处理每张：
      * "template_idx": 修改既有 sample slide 字段
      * "add_layout":   add_slide_with_layout（仅安全 layout，如空白页5）
      * "duplicate":    复制既有 sample slide → 改字段
  - 最后删除未使用的 sample slide + 重排顺序

slide_structure.json 格式 v0.1：
{
  "slides": [
    {
      "slide_num": 1,
      "source": {"type": "template_idx", "idx": 1},
      "placeholders": {"12": "...", "13": "..."},
      "free_shapes": {"文本框 23": "春季"}
    },
    {
      "slide_num": 5,
      "source": {"type": "add_layout", "layout": "空白页5"},
      "placeholders": {"10": "01.", "11": "分解质因数的方法"}
    },
    {
      "slide_num": 2,
      "source": {"type": "duplicate", "from_idx": 2},
      "placeholders": {...},
      "free_shapes": {...}
    }
  ]
}

字段索引（template_idx）从 1 开始（用户友好）。
"""

import argparse
import copy
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.oxml.ns import qn
from lxml import etree

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ──────────────────────────────────────────────────
# 字体规则
# ──────────────────────────────────────────────────

def _set_run_lang(run, font_size=None):
    rpr = run._r.get_or_add_rPr()
    rpr.set("lang", "zh-CN")
    ea = rpr.find(f"{{{A_NS}}}ea")
    if ea is None:
        ea = etree.SubElement(rpr, f"{{{A_NS}}}ea")
    ea.set("typeface", "微软雅黑")
    if font_size:
        run.font.size = Pt(font_size)


# ──────────────────────────────────────────────────
# Shape 操作
# ──────────────────────────────────────────────────

def _write_text_frame(shape, text, font_size=None):
    """把文本写入 shape.text_frame（清空原内容）。"""
    if not shape.has_text_frame:
        print(f"  ⚠ shape {shape.name!r} 没 text_frame，跳过", file=sys.stderr)
        return
    tf = shape.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = para
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        _set_run_lang(run, font_size)


def _find_placeholder(slide, ph_idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            return ph
    return None


def _ensure_placeholder(slide, ph_idx):
    """缺占位符时从 layout 复制 sp 补齐。"""
    p = _find_placeholder(slide, ph_idx)
    if p is not None:
        return p
    layout = slide.slide_layout
    src = None
    for lph in layout.placeholders:
        if lph.placeholder_format.idx == ph_idx:
            src = lph._element
            break
    if src is None:
        return None
    new_sp = copy.deepcopy(src)
    slide.shapes._spTree.append(new_sp)
    return _find_placeholder(slide, ph_idx)


def _find_free_shape(slide, shape_name):
    """按 name 查 free（非 placeholder）shape。"""
    for shape in slide.shapes:
        if shape.name == shape_name and not shape.is_placeholder:
            return shape
    return None


def _write_placeholders(slide, placeholders_dict):
    """填多个 placeholder。"""
    for ph_idx_str, text in placeholders_dict.items():
        ph_idx = int(ph_idx_str)
        target = _ensure_placeholder(slide, ph_idx)
        if target is None:
            print(f"    ⚠ ph[{ph_idx}] 在此 slide 找不到，跳过", file=sys.stderr)
            continue
        _write_text_frame(target, text)


def _resize_placeholders(slide, resize_dict):
    """调整 placeholder 的大小/位置。

    格式：{"16": {"width": 12.0, "left": 0.65, "top": 1.0, "height": 0.5}}
    （单位英寸，未提供的字段保持不变）
    """
    for ph_idx_str, dims in resize_dict.items():
        ph_idx = int(ph_idx_str)
        target = _ensure_placeholder(slide, ph_idx)
        if target is None:
            continue
        if "left" in dims:
            target.left = Inches(dims["left"])
        if "top" in dims:
            target.top = Inches(dims["top"])
        if "width" in dims:
            target.width = Inches(dims["width"])
        if "height" in dims:
            target.height = Inches(dims["height"])


def _write_free_shapes(slide, free_shapes_dict):
    """按 name 修改 free shape（非 placeholder）的文字。

    value 可以是：
      - str: 仅文字
      - dict: {"text": str, "font_size": int}
    """
    for shape_name, value in free_shapes_dict.items():
        target = _find_free_shape(slide, shape_name)
        if target is None:
            print(f"    ⚠ free shape {shape_name!r} 找不到，跳过", file=sys.stderr)
            continue
        if isinstance(value, dict):
            text = value["text"]
            font_size = value.get("font_size")
        else:
            text = value
            font_size = None
        _write_text_frame(target, text, font_size=font_size)


def _add_text_box(slide, text, left_inch, top_inch, width_inch, height_inch, font_size=20):
    """在 slide 上新增一个自由文本框（用于解析/思考题/笔记正文等）。"""
    tb = slide.shapes.add_textbox(
        Inches(left_inch), Inches(top_inch),
        Inches(width_inch), Inches(height_inch)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = para if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        _set_run_lang(run, font_size)
    return tb


def _add_text_boxes(slide, boxes_list):
    """根据 list of dicts 添加多个文本框。"""
    for box in boxes_list:
        _add_text_box(
            slide,
            text=box["text"],
            left_inch=box.get("left", 1.0),
            top_inch=box.get("top", 1.0),
            width_inch=box.get("width", 10.0),
            height_inch=box.get("height", 2.0),
            font_size=box.get("font_size", 20),
        )


def _add_image(slide, image_path, left_inch, top_inch, width_inch=None, height_inch=None):
    """在 slide 上插入图片。可只给 width 或 height（按比例自动算另一边）。"""
    from pptx.util import Inches
    kwargs = {}
    if width_inch is not None:
        kwargs["width"] = Inches(width_inch)
    if height_inch is not None:
        kwargs["height"] = Inches(height_inch)
    slide.shapes.add_picture(
        image_path,
        Inches(left_inch), Inches(top_inch),
        **kwargs,
    )


def _add_images(slide, images_list):
    """根据 list of dicts 添加多张图片。"""
    for img in images_list:
        _add_image(
            slide,
            image_path=img["image_path"],
            left_inch=img.get("left", 1.0),
            top_inch=img.get("top", 1.0),
            width_inch=img.get("width"),
            height_inch=img.get("height"),
        )


def _add_table(slide, data, left_inch, top_inch, width_inch, height_inch,
               font_size=14, header_fill=None):
    """在 slide 上插入表格。data 是 list of list of str（含表头行）。"""
    from pptx.dml.color import RGBColor
    rows = len(data)
    cols = max(len(r) for r in data) if data else 0
    if rows == 0 or cols == 0:
        return
    shape = slide.shapes.add_table(
        rows, cols,
        Inches(left_inch), Inches(top_inch),
        Inches(width_inch), Inches(height_inch),
    )
    table = shape.table
    for r, row in enumerate(data):
        for c in range(cols):
            cell = table.cell(r, c)
            text = row[c] if c < len(row) else ""
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = str(text)
            _set_run_lang(run, font_size)
            if r == 0 and header_fill:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(header_fill.lstrip("#"))
                run.font.bold = True


def _add_tables(slide, tables_list):
    """根据 list of dicts 添加多个表格。"""
    for t in tables_list:
        _add_table(
            slide,
            data=t["data"],
            left_inch=t.get("left", 1.0),
            top_inch=t.get("top", 3.0),
            width_inch=t.get("width", 11.0),
            height_inch=t.get("height", 2.0),
            font_size=t.get("font_size", 14),
            header_fill=t.get("header_fill"),
        )


# ──────────────────────────────────────────────────
# Slide 操作
# ──────────────────────────────────────────────────

def _get_slide_by_idx(prs, idx_1based):
    return prs.slides[idx_1based - 1]


def _add_slide_with_layout(prs, layout_name):
    for L in prs.slide_layouts:
        if L.name == layout_name:
            return prs.slides.add_slide(L)
    raise ValueError(f"模板里没有版式 {layout_name!r}")


def _duplicate_slide(prs, source_idx_1based):
    """通过 lxml 复制既有 slide（保留所有 shape 含装饰）→ 加到 prs 末尾，返回新 slide。"""
    source_slide = _get_slide_by_idx(prs, source_idx_1based)
    # 用 source 的 layout 新建一张 slide
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    # 清除新 slide 上 layout 自动复制过来的占位符（保留 spTree 的根框架）
    # 然后从 source 复制所有 shape 过来
    new_sp_tree = new_slide.shapes._spTree
    src_sp_tree = source_slide.shapes._spTree

    # 删除 new_slide 现有的所有 sp/grpSp/pic（保留 nvGrpSpPr 和 grpSpPr 这种 group container）
    for child in list(new_sp_tree):
        tag = etree.QName(child).localname
        if tag in ("sp", "grpSp", "pic", "graphicFrame", "cxnSp"):
            new_sp_tree.remove(child)

    # 从 source 复制所有 sp 等到 new
    for child in src_sp_tree:
        tag = etree.QName(child).localname
        if tag in ("sp", "grpSp", "pic", "graphicFrame", "cxnSp"):
            new_sp_tree.append(copy.deepcopy(child))

    return new_slide


def _delete_slide_by_idx(prs, idx_1based):
    sld_id_lst = prs.slides._sldIdLst
    sld_id_elem = list(sld_id_lst)[idx_1based - 1]
    rId = sld_id_elem.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sld_id_lst.remove(sld_id_elem)


def _move_slide(prs, current_idx_1based, target_idx_1based):
    sld_id_lst = prs.slides._sldIdLst
    elem = list(sld_id_lst)[current_idx_1based - 1]
    sld_id_lst.remove(elem)
    sld_id_lst.insert(target_idx_1based - 1, elem)


# ──────────────────────────────────────────────────
# 主流程
# ──────────────────────────────────────────────────

def render(template_path: Path, structure_path: Path, output_path: Path):
    structure = json.loads(structure_path.read_text(encoding="utf-8"))
    slides_data = structure.get("slides", [])
    if not slides_data:
        sys.exit("❌ slide_structure.json 里没有 slides")

    prs = Presentation(str(template_path))
    original_count = len(prs.slides)
    print(f"加载模板：{template_path}")
    print(f"  原 slide 数：{original_count}")

    # 跟踪每个 slide_num 对应的 slide 对象（用 element 比对而不是 idx）
    slide_num_to_element = {}

    # ⚠️ 阶段顺序很重要：先 duplicate（基于原始 sample），再 modify。
    # 否则 modify 后 duplicate 会复制到已被修改过的内容，导致重复叠加。

    # 阶段 1（前）：处理 duplicate 类（追加在末尾，基于原始 sample）
    print("\n阶段 1：复制 sample slide（在 modify 之前，保证拿到原始 sample）")
    for sd in slides_data:
        source = sd["source"]
        if source["type"] != "duplicate":
            continue
        from_idx = source["from_idx"]
        new_slide = _duplicate_slide(prs, from_idx)
        _write_placeholders(new_slide, sd.get("placeholders", {}))
        _write_free_shapes(new_slide, sd.get("free_shapes", {}))
        _add_text_boxes(new_slide, sd.get("add_text_boxes", []))
        _add_images(new_slide, sd.get("add_images", []))
        _add_tables(new_slide, sd.get("add_tables", []))
        new_idx = len(prs.slides)
        sld_id = list(prs.slides._sldIdLst)[new_idx - 1]
        slide_num_to_element[sd["slide_num"]] = sld_id
        print(f"  slide_num {sd['slide_num']:2d} ← duplicated from idx {from_idx} (new idx={new_idx})")

    # 阶段 2：处理 template_idx 类（修改既有 sample slide 字段）
    used_template_idx = set()
    print("\n阶段 2：修改既有 sample slide")
    for sd in slides_data:
        source = sd["source"]
        if source["type"] != "template_idx":
            continue
        idx = source["idx"]
        slide = _get_slide_by_idx(prs, idx)
        _write_placeholders(slide, sd.get("placeholders", {}))
        _write_free_shapes(slide, sd.get("free_shapes", {}))
        _add_text_boxes(slide, sd.get("add_text_boxes", []))
        _add_images(slide, sd.get("add_images", []))
        _add_tables(slide, sd.get("add_tables", []))
        used_template_idx.add(idx)
        # 用 slide id 跟踪
        sld_id = list(prs.slides._sldIdLst)[idx - 1]
        slide_num_to_element[sd["slide_num"]] = sld_id
        print(f"  slide_num {sd['slide_num']:2d} ← template idx {idx}")

    # 阶段 3：处理 add_layout 类（追加在末尾）
    print("\n阶段 3：add_slide 新增（仅安全 layout）")
    for sd in slides_data:
        source = sd["source"]
        if source["type"] != "add_layout":
            continue
        layout_name = source["layout"]
        new_slide = _add_slide_with_layout(prs, layout_name)
        _write_placeholders(new_slide, sd.get("placeholders", {}))
        _write_free_shapes(new_slide, sd.get("free_shapes", {}))
        _add_text_boxes(new_slide, sd.get("add_text_boxes", []))
        _add_images(new_slide, sd.get("add_images", []))
        _add_tables(new_slide, sd.get("add_tables", []))
        new_idx = len(prs.slides)
        sld_id = list(prs.slides._sldIdLst)[new_idx - 1]
        slide_num_to_element[sd["slide_num"]] = sld_id
        print(f"  slide_num {sd['slide_num']:2d} ← add_layout({layout_name!r}) (new idx={new_idx})")

    # 阶段 4：删除未使用的模板 sample slide
    print("\n阶段 4：删除未使用的 sample slide")
    unused = [i for i in range(1, original_count + 1) if i not in used_template_idx]
    print(f"  未使用 sample idx: {unused}")
    # 从后往前删，避免 idx 偏移
    for idx in sorted(unused, reverse=True):
        _delete_slide_by_idx(prs, idx)
        print(f"  删除 idx {idx}")

    # 阶段 5：按 slide_num 严格重排（用 sldId element 直接重排，不再靠文本定位）
    print("\n阶段 5：按 slide_num 重排")
    sld_id_lst = prs.slides._sldIdLst
    sorted_sds = sorted(slides_data, key=lambda x: x["slide_num"])
    # 先把所有目标 element 从 sld_id_lst 摘出来
    target_elements = []
    for sd in sorted_sds:
        sld_id = slide_num_to_element.get(sd["slide_num"])
        if sld_id is None:
            print(f"  ⚠ slide_num {sd['slide_num']} 没建立 element 追踪，跳过", file=sys.stderr)
            continue
        target_elements.append((sd["slide_num"], sld_id))

    # 全部 detach
    for _, sld_id in target_elements:
        if sld_id in sld_id_lst:
            sld_id_lst.remove(sld_id)
    # 按目标顺序 append（注意：sld_id_lst 此时应该是空的——除非有遗漏的未处理 slide）
    remaining = list(sld_id_lst)
    for r in remaining:
        sld_id_lst.remove(r)
        print(f"  ⚠ 有未追踪 slide 被强制删除", file=sys.stderr)
    for slide_num, sld_id in target_elements:
        sld_id_lst.append(sld_id)
        print(f"  位置 → slide_num {slide_num}")

    # 阶段 6：保存
    out_path = Path(output_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"\n✅ 渲染完成：{out_path}")
    print(f"   总 slide 数：{len(prs.slides)}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--slide-structure", required=True)
    parser.add_argument("--template", required=True)
    parser.add_argument("-o", "--output", required=True)
    args = parser.parse_args()
    render(Path(args.template), Path(args.slide_structure), Path(args.output))


if __name__ == "__main__":
    main()
